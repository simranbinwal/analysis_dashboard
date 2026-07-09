from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from typing import Any

import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import streamlit as st
from plotly.subplots import make_subplots
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
from sklearn.cluster import KMeans
from sklearn.decomposition import PCA
from sklearn.ensemble import IsolationForest, RandomForestRegressor
from sklearn.impute import SimpleImputer
from sklearn.linear_model import LinearRegression
from sklearn.metrics import r2_score
from sklearn.pipeline import Pipeline
from sklearn.preprocessing import StandardScaler
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


st.set_page_config(
    page_title="AI Data Cleaning and Visualization System",
    page_icon="AD",
    layout="wide",
    initial_sidebar_state="expanded",
)


@dataclass
class AnalysisResult:
    raw_df: pd.DataFrame
    clean_df: pd.DataFrame
    numeric_cols: list[str]
    categorical_cols: list[str]
    date_cols: list[str]
    missing_values: int
    duplicates_removed: int
    outlier_count: int
    outlier_columns: dict[str, int]
    quality: dict[str, int]
    insights: list[str]
    forecast: pd.DataFrame | None
    forecast_target: str | None
    forecast_score: float | None
    segments: pd.DataFrame | None
    segment_profile: pd.DataFrame | None


def inject_css() -> None:
    st.markdown(
        """
        <style>
        :root {
            --bg: #07111f;
            --panel: rgba(12, 23, 42, 0.78);
            --panel-2: rgba(16, 33, 59, 0.82);
            --text: #eaf2ff;
            --muted: #9eb2cb;
            --cyan: #38d5ff;
            --green: #2df5a1;
            --amber: #f7c65b;
            --pink: #ff4da6;
            --line: rgba(123, 213, 255, 0.18);
        }

        .stApp {
            background:
                radial-gradient(circle at 18% 10%, rgba(56, 213, 255, 0.18), transparent 28%),
                radial-gradient(circle at 84% 18%, rgba(255, 77, 166, 0.14), transparent 30%),
                linear-gradient(135deg, #06101d 0%, #0b1729 42%, #09111f 100%);
            color: var(--text);
        }

        [data-testid="stSidebar"] {
            background: linear-gradient(180deg, rgba(7, 17, 31, 0.98), rgba(12, 24, 43, 0.96));
            border-right: 1px solid var(--line);
        }

        [data-testid="stSidebar"] * {
            color: var(--text);
        }

        .hero {
            border: 1px solid var(--line);
            border-radius: 18px;
            padding: 28px;
            background:
                linear-gradient(135deg, rgba(56, 213, 255, 0.14), rgba(45, 245, 161, 0.08)),
                rgba(8, 18, 34, 0.84);
            box-shadow: 0 22px 80px rgba(0, 0, 0, 0.28);
            margin-bottom: 18px;
        }

        .hero h1 {
            color: var(--text);
            font-size: 2.35rem;
            line-height: 1.12;
            margin: 0 0 8px;
            letter-spacing: 0;
        }

        .hero p {
            color: var(--muted);
            font-size: 1rem;
            margin: 0;
            max-width: 920px;
        }

        .metric-card {
            min-height: 128px;
            padding: 18px;
            border-radius: 16px;
            border: 1px solid var(--line);
            background: linear-gradient(180deg, var(--panel), rgba(8, 18, 34, 0.68));
            box-shadow: inset 0 1px 0 rgba(255, 255, 255, 0.06), 0 18px 44px rgba(0, 0, 0, 0.22);
        }

        .metric-card span {
            color: var(--muted);
            display: block;
            font-size: 0.82rem;
            font-weight: 700;
            text-transform: uppercase;
            letter-spacing: 0;
        }

        .metric-card strong {
            color: var(--text);
            display: block;
            font-size: 2.05rem;
            margin-top: 10px;
        }

        .metric-card small {
            color: var(--green);
        }

        .glass-panel {
            padding: 18px;
            border-radius: 16px;
            border: 1px solid var(--line);
            background: var(--panel);
            box-shadow: 0 20px 64px rgba(0, 0, 0, 0.22);
        }

        .section-title {
            color: var(--text);
            font-size: 1.16rem;
            font-weight: 800;
            margin-bottom: 8px;
        }

        .insight {
            padding: 12px 14px;
            margin: 8px 0;
            border-radius: 12px;
            color: var(--text);
            background: rgba(56, 213, 255, 0.08);
            border-left: 4px solid var(--cyan);
        }

        .chat-user {
            padding: 12px 14px;
            margin: 8px 0 8px auto;
            max-width: 82%;
            border-radius: 14px;
            background: rgba(45, 245, 161, 0.13);
            border: 1px solid rgba(45, 245, 161, 0.24);
        }

        .chat-ai {
            padding: 12px 14px;
            margin: 8px 0;
            max-width: 86%;
            border-radius: 14px;
            background: rgba(56, 213, 255, 0.12);
            border: 1px solid rgba(56, 213, 255, 0.24);
        }

        .stTabs [data-baseweb="tab-list"] {
            gap: 10px;
        }

        .stTabs [data-baseweb="tab"] {
            border-radius: 999px;
            border: 1px solid var(--line);
            background: rgba(12, 23, 42, 0.76);
            padding: 10px 18px;
        }

        .stDataFrame, .stTable {
            border: 1px solid var(--line);
            border-radius: 14px;
            overflow: hidden;
        }

        div[data-testid="stDownloadButton"] button,
        div[data-testid="stButton"] button {
            border: 1px solid rgba(56, 213, 255, 0.36);
            border-radius: 12px;
            background: linear-gradient(135deg, rgba(56, 213, 255, 0.88), rgba(45, 245, 161, 0.78));
            color: #06101d;
            font-weight: 900;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )


def sample_dataset() -> pd.DataFrame:
    dates = pd.date_range("2025-01-01", periods=84, freq="D")
    rng = np.random.default_rng(42)
    regions = np.array(["North", "South", "East", "West"])
    products = np.array(["Laptop", "Phone", "Tablet", "Accessory"])
    customers = np.array(
        [
            "Aster Retail",
            "Bright Mart",
            "Cascade Inc",
            "Dwell Co",
            "Elm Foods",
            "Focal Shop",
            "Grove Stores",
            "Halo Market",
        ]
    )

    df = pd.DataFrame(
        {
            "Date": dates,
            "Region": rng.choice(regions, len(dates)),
            "Customer": rng.choice(customers, len(dates)),
            "Product": rng.choice(products, len(dates)),
            "Sales": np.linspace(4200, 18800, len(dates)) + rng.normal(0, 1300, len(dates)),
            "Quantity": rng.integers(4, 36, len(dates)),
            "Discount": rng.uniform(0.03, 0.24, len(dates)),
            "Satisfaction": rng.uniform(6.8, 9.7, len(dates)),
        }
    )
    df.loc[7, "Sales"] = np.nan
    df.loc[13, "Satisfaction"] = np.nan
    df.loc[31, "Sales"] = 52000
    df = pd.concat([df, df.iloc[[10]]], ignore_index=True)
    return df.round({"Sales": 2, "Discount": 2, "Satisfaction": 1})


def read_uploaded_file(uploaded_file: Any) -> pd.DataFrame:
    if uploaded_file.name.lower().endswith(".csv"):
        return pd.read_csv(uploaded_file)
    if uploaded_file.name.lower().endswith((".xlsx", ".xls")):
        return pd.read_excel(uploaded_file)
    raise ValueError("Please upload a CSV or Excel file.")


def normalize_missing(df: pd.DataFrame) -> pd.DataFrame:
    missing_tokens = ["", " ", "na", "n/a", "none", "null", "nan", "-", "--"]
    return df.replace(missing_tokens, np.nan)


def infer_columns(df: pd.DataFrame) -> tuple[list[str], list[str], list[str]]:
    numeric_cols: list[str] = []
    date_cols: list[str] = []

    for column in df.columns:
        series = df[column].dropna()
        if series.empty:
            continue

        numeric_ratio = pd.to_numeric(series, errors="coerce").notna().mean()
        date_ratio = pd.to_datetime(series, errors="coerce").notna().mean()

        if numeric_ratio >= 0.75:
            numeric_cols.append(column)
        elif date_ratio >= 0.75:
            date_cols.append(column)

    categorical_cols = [column for column in df.columns if column not in numeric_cols + date_cols]
    return numeric_cols, categorical_cols, date_cols


def detect_and_cap_outliers(df: pd.DataFrame, numeric_cols: list[str]) -> tuple[pd.DataFrame, int, dict[str, int]]:
    capped = df.copy()
    outlier_columns: dict[str, int] = {}
    total = 0

    for column in numeric_cols:
        values = capped[column].dropna()
        if values.size < 4:
            continue

        q1 = values.quantile(0.25)
        q3 = values.quantile(0.75)
        iqr = q3 - q1
        if iqr == 0:
            continue

        lower = q1 - 1.5 * iqr
        upper = q3 + 1.5 * iqr
        mask = (capped[column] < lower) | (capped[column] > upper)
        count = int(mask.sum())
        if count:
            capped.loc[capped[column] < lower, column] = lower
            capped.loc[capped[column] > upper, column] = upper
            outlier_columns[column] = count
            total += count

    return capped, total, outlier_columns


def clean_dataset(raw_df: pd.DataFrame) -> tuple[pd.DataFrame, dict[str, Any]]:
    raw = normalize_missing(raw_df.copy())
    missing_values = int(raw.isna().sum().sum())
    duplicate_count = int(raw.duplicated().sum())
    deduped = raw.drop_duplicates().reset_index(drop=True)
    numeric_cols, categorical_cols, date_cols = infer_columns(deduped)

    cleaned = deduped.copy()
    for column in numeric_cols:
        cleaned[column] = pd.to_numeric(cleaned[column], errors="coerce")

    for column in date_cols:
        cleaned[column] = pd.to_datetime(cleaned[column], errors="coerce")

    if numeric_cols:
        numeric_imputer = SimpleImputer(strategy="median")
        cleaned[numeric_cols] = numeric_imputer.fit_transform(cleaned[numeric_cols])

    for column in categorical_cols:
        mode = cleaned[column].mode(dropna=True)
        fill_value = mode.iloc[0] if not mode.empty else "Unknown"
        cleaned[column] = cleaned[column].fillna(fill_value).astype(str)

    for column in date_cols:
        mode = cleaned[column].mode(dropna=True)
        fill_value = mode.iloc[0] if not mode.empty else pd.Timestamp.today().normalize()
        cleaned[column] = cleaned[column].fillna(fill_value)

    cleaned, outlier_count, outlier_columns = detect_and_cap_outliers(cleaned, numeric_cols)
    metadata = {
        "missing_values": missing_values,
        "duplicates_removed": duplicate_count,
        "outlier_count": outlier_count,
        "outlier_columns": outlier_columns,
        "numeric_cols": numeric_cols,
        "categorical_cols": categorical_cols,
        "date_cols": date_cols,
    }
    return cleaned, metadata


def data_quality_score(raw_df: pd.DataFrame, metadata: dict[str, Any]) -> dict[str, int]:
    cells = max(1, raw_df.shape[0] * raw_df.shape[1])
    completeness = 1 - metadata["missing_values"] / cells
    uniqueness = 1 - metadata["duplicates_removed"] / max(1, raw_df.shape[0])
    outlier_health = 1 - min(metadata["outlier_count"] / max(1, raw_df.shape[0]), 1) * 0.5

    consistency_checks = 0
    consistency_pass = 0
    normalized = normalize_missing(raw_df.copy())
    for column in metadata["numeric_cols"]:
        values = normalized[column].dropna()
        consistency_checks += len(values)
        consistency_pass += int(pd.to_numeric(values, errors="coerce").notna().sum())
    for column in metadata["date_cols"]:
        values = normalized[column].dropna()
        consistency_checks += len(values)
        consistency_pass += int(pd.to_datetime(values, errors="coerce").notna().sum())
    for column in metadata["categorical_cols"]:
        consistency_checks += int(normalized[column].dropna().shape[0])
        consistency_pass += int(normalized[column].dropna().shape[0])

    consistency = consistency_pass / consistency_checks if consistency_checks else 1
    score = completeness * 0.35 + uniqueness * 0.25 + consistency * 0.2 + outlier_health * 0.2
    return {
        "score": int(round(score * 100)),
        "completeness": int(round(completeness * 100)),
        "uniqueness": int(round(uniqueness * 100)),
        "consistency": int(round(consistency * 100)),
        "outlier_health": int(round(outlier_health * 100)),
    }


def choose_target(numeric_cols: list[str]) -> str | None:
    if not numeric_cols:
        return None
    priority = ["sales", "revenue", "amount", "profit", "value", "price", "target"]
    for term in priority:
        for column in numeric_cols:
            if term in column.lower():
                return column
    return numeric_cols[0]


def build_forecast(clean_df: pd.DataFrame, numeric_cols: list[str], date_cols: list[str]) -> tuple[pd.DataFrame | None, str | None, float | None]:
    target = choose_target(numeric_cols)
    if target is None or len(clean_df) < 5:
        return None, None, None

    if date_cols:
        date_col = date_cols[0]
        series = (
            clean_df[[date_col, target]]
            .dropna()
            .groupby(pd.Grouper(key=date_col, freq="D"))[target]
            .sum()
            .reset_index()
            .sort_values(date_col)
        )
        series = series.rename(columns={date_col: "Period", target: "Actual"})
    else:
        series = pd.DataFrame({"Period": np.arange(1, len(clean_df) + 1), "Actual": clean_df[target].values})

    series = series.dropna().reset_index(drop=True)
    if len(series) < 5:
        return None, target, None

    x = np.arange(len(series)).reshape(-1, 1)
    y = series["Actual"].to_numpy()
    model = LinearRegression()
    model.fit(x, y)
    fitted = model.predict(x)
    score = float(r2_score(y, fitted)) if len(np.unique(y)) > 1 else 1.0
    future_x = np.arange(len(series), len(series) + 7).reshape(-1, 1)
    future_y = model.predict(future_x)

    if date_cols:
        last_date = pd.to_datetime(series["Period"].iloc[-1])
        future_period = pd.date_range(last_date + pd.Timedelta(days=1), periods=7, freq="D")
    else:
        future_period = np.arange(len(series) + 1, len(series) + 8)

    forecast = pd.concat(
        [
            pd.DataFrame({"Period": series["Period"], "Actual": series["Actual"], "Forecast": fitted, "Type": "Historical"}),
            pd.DataFrame({"Period": future_period, "Actual": np.nan, "Forecast": future_y, "Type": "Forecast"}),
        ],
        ignore_index=True,
    )
    forecast["Forecast"] = forecast["Forecast"].clip(lower=0)
    return forecast, target, score


def build_segments(clean_df: pd.DataFrame, numeric_cols: list[str]) -> tuple[pd.DataFrame | None, pd.DataFrame | None]:
    if len(numeric_cols) < 2 or len(clean_df) < 6:
        return None, None

    feature_cols = numeric_cols[: min(5, len(numeric_cols))]
    features = clean_df[feature_cols].dropna()
    if len(features) < 6:
        return None, None

    cluster_count = min(4, max(2, int(np.sqrt(len(features) / 2))))
    pipeline = Pipeline([("scaler", StandardScaler()), ("kmeans", KMeans(n_clusters=cluster_count, random_state=42, n_init=10))])
    labels = pipeline.fit_predict(features)

    pca = PCA(n_components=2, random_state=42)
    projected = pca.fit_transform(StandardScaler().fit_transform(features))
    segments = clean_df.loc[features.index].copy()
    segments["Segment"] = [f"Segment {label + 1}" for label in labels]
    segments["PCA 1"] = projected[:, 0]
    segments["PCA 2"] = projected[:, 1]

    target = choose_target(numeric_cols) or feature_cols[0]
    profile = (
        segments.groupby("Segment")
        .agg(Records=("Segment", "size"), Average_Target=(target, "mean"))
        .reset_index()
        .sort_values("Records", ascending=False)
    )
    return segments, profile


def anomaly_signal(clean_df: pd.DataFrame, numeric_cols: list[str]) -> pd.DataFrame | None:
    if len(numeric_cols) < 2 or len(clean_df) < 10:
        return None
    features = clean_df[numeric_cols].dropna()
    model = IsolationForest(contamination="auto", random_state=42)
    labels = model.fit_predict(features)
    output = clean_df.loc[features.index].copy()
    output["Anomaly"] = np.where(labels == -1, "Anomaly", "Normal")
    return output


def predictive_importance(clean_df: pd.DataFrame, numeric_cols: list[str]) -> pd.DataFrame | None:
    target = choose_target(numeric_cols)
    if target is None or len(numeric_cols) < 2:
        return None
    features = [column for column in numeric_cols if column != target]
    if not features:
        return None

    x = clean_df[features].dropna()
    y = clean_df.loc[x.index, target]
    if len(x) < 8:
        return None

    model = RandomForestRegressor(n_estimators=120, random_state=42)
    model.fit(x, y)
    return pd.DataFrame({"Feature": features, "Importance": model.feature_importances_}).sort_values("Importance", ascending=False)


def generate_insights(
    clean_df: pd.DataFrame,
    metadata: dict[str, Any],
    quality: dict[str, int],
    forecast: pd.DataFrame | None,
    forecast_target: str | None,
    segment_profile: pd.DataFrame | None,
) -> list[str]:
    insights = [
        f"{len(clean_df):,} clean records are ready after removing {metadata['duplicates_removed']:,} duplicate rows.",
        f"{metadata['missing_values']:,} missing values were handled with median, mode, or date imputation.",
        f"Data quality score is {quality['score']}/100 with {quality['completeness']}% completeness and {quality['uniqueness']}% uniqueness.",
    ]

    if metadata["outlier_count"]:
        top_outlier = max(metadata["outlier_columns"].items(), key=lambda item: item[1])[0]
        insights.append(f"{metadata['outlier_count']:,} outliers were detected and capped; {top_outlier} had the strongest outlier signal.")
    else:
        insights.append("No major IQR outliers were detected in numeric columns.")

    if forecast is not None and forecast_target is not None:
        next_value = forecast.loc[forecast["Type"] == "Forecast", "Forecast"].iloc[0]
        insights.append(f"Next forecasted {forecast_target} value is approximately {next_value:,.2f}.")

    if segment_profile is not None and not segment_profile.empty:
        largest = segment_profile.iloc[0]
        insights.append(f"{largest['Segment']} is the largest customer segment with {int(largest['Records'])} records.")

    return insights


def analyze(raw_df: pd.DataFrame) -> AnalysisResult:
    clean_df, metadata = clean_dataset(raw_df)
    quality = data_quality_score(raw_df, metadata)
    forecast, forecast_target, forecast_score = build_forecast(clean_df, metadata["numeric_cols"], metadata["date_cols"])
    segments, segment_profile = build_segments(clean_df, metadata["numeric_cols"])
    insights = generate_insights(clean_df, metadata, quality, forecast, forecast_target, segment_profile)

    return AnalysisResult(
        raw_df=raw_df,
        clean_df=clean_df,
        numeric_cols=metadata["numeric_cols"],
        categorical_cols=metadata["categorical_cols"],
        date_cols=metadata["date_cols"],
        missing_values=metadata["missing_values"],
        duplicates_removed=metadata["duplicates_removed"],
        outlier_count=metadata["outlier_count"],
        outlier_columns=metadata["outlier_columns"],
        quality=quality,
        insights=insights,
        forecast=forecast,
        forecast_target=forecast_target,
        forecast_score=forecast_score,
        segments=segments,
        segment_profile=segment_profile,
    )


def metric_card(label: str, value: str, note: str = "") -> None:
    st.markdown(
        f"""
        <div class="metric-card">
            <span>{label}</span>
            <strong>{value}</strong>
            <small>{note}</small>
        </div>
        """,
        unsafe_allow_html=True,
    )


def show_hero() -> None:
    st.markdown(
        """
        <div class="hero">
            <h1>AI-Powered Automated Data Cleaning and Visualization System</h1>
            <p>
                Upload a dataset and the system cleans missing values, removes duplicates, detects outliers,
                scores data quality, builds interactive dashboards, forecasts sales-like metrics, segments customers,
                and exports executive-ready reports.
            </p>
        </div>
        """,
        unsafe_allow_html=True,
    )


def quality_gauge(score: int) -> go.Figure:
    fig = go.Figure(
        go.Indicator(
            mode="gauge+number",
            value=score,
            number={"suffix": "/100", "font": {"color": "#eaf2ff", "size": 34}},
            gauge={
                "axis": {"range": [0, 100], "tickcolor": "#9eb2cb"},
                "bar": {"color": "#2df5a1"},
                "bgcolor": "rgba(12,23,42,0.5)",
                "borderwidth": 1,
                "bordercolor": "rgba(123,213,255,0.25)",
                "steps": [
                    {"range": [0, 60], "color": "rgba(255,77,166,0.28)"},
                    {"range": [60, 80], "color": "rgba(247,198,91,0.28)"},
                    {"range": [80, 100], "color": "rgba(45,245,161,0.24)"},
                ],
            },
        )
    )
    return style_plot(fig, height=280)


def style_plot(fig: go.Figure, height: int = 420) -> go.Figure:
    fig.update_layout(
        height=height,
        paper_bgcolor="rgba(0,0,0,0)",
        plot_bgcolor="rgba(7,17,31,0.28)",
        font={"color": "#eaf2ff"},
        margin=dict(l=24, r=24, t=56, b=42),
        legend=dict(bgcolor="rgba(0,0,0,0)"),
    )
    fig.update_xaxes(gridcolor="rgba(123,213,255,0.12)", zerolinecolor="rgba(123,213,255,0.14)")
    fig.update_yaxes(gridcolor="rgba(123,213,255,0.12)", zerolinecolor="rgba(123,213,255,0.14)")
    return fig


def dashboard_tab(result: AnalysisResult) -> None:
    target = choose_target(result.numeric_cols)
    left, right = st.columns([1.35, 0.9])

    with left:
        st.markdown('<div class="section-title">Auto Dashboard Builder</div>', unsafe_allow_html=True)
        chart_type = st.segmented_control(
            "Chart mode",
            ["Distribution", "Category Summary", "Correlation", "Timeline"],
            default="Distribution",
            label_visibility="collapsed",
        )

        if chart_type == "Distribution" and result.numeric_cols:
            column = st.selectbox("Numeric field", result.numeric_cols, key="dist_col")
            fig = px.histogram(result.clean_df, x=column, nbins=18, color_discrete_sequence=["#38d5ff"])
            st.plotly_chart(style_plot(fig), use_container_width=True)
        elif chart_type == "Category Summary" and result.categorical_cols and target:
            category = st.selectbox("Category field", result.categorical_cols, key="cat_col")
            grouped = result.clean_df.groupby(category)[target].sum().sort_values(ascending=False).head(12).reset_index()
            fig = px.bar(grouped, x=category, y=target, color=target, color_continuous_scale="Tealgrn")
            st.plotly_chart(style_plot(fig), use_container_width=True)
        elif chart_type == "Correlation" and len(result.numeric_cols) >= 2:
            corr = result.clean_df[result.numeric_cols].corr()
            fig = px.imshow(corr, text_auto=True, color_continuous_scale="Bluered", aspect="auto")
            st.plotly_chart(style_plot(fig), use_container_width=True)
        elif chart_type == "Timeline" and result.date_cols and target:
            date_col = result.date_cols[0]
            timeline = result.clean_df.groupby(pd.Grouper(key=date_col, freq="D"))[target].sum().reset_index()
            fig = px.line(timeline, x=date_col, y=target, markers=True, color_discrete_sequence=["#2df5a1"])
            st.plotly_chart(style_plot(fig), use_container_width=True)
        else:
            st.info("This chart needs compatible numeric, category, or date columns.")

    with right:
        st.markdown('<div class="section-title">Data Quality Score</div>', unsafe_allow_html=True)
        st.plotly_chart(quality_gauge(result.quality["score"]), use_container_width=True)
        components = pd.DataFrame(
            {
                "Component": ["Completeness", "Uniqueness", "Consistency", "Outlier Health"],
                "Score": [
                    result.quality["completeness"],
                    result.quality["uniqueness"],
                    result.quality["consistency"],
                    result.quality["outlier_health"],
                ],
            }
        )
        fig = px.bar(components, x="Component", y="Score", color="Score", color_continuous_scale="Mint")
        st.plotly_chart(style_plot(fig, height=260), use_container_width=True)


def cleaning_tab(result: AnalysisResult) -> None:
    col1, col2 = st.columns([0.9, 1.1])
    with col1:
        st.markdown('<div class="section-title">Automated Cleaning Summary</div>', unsafe_allow_html=True)
        summary = pd.DataFrame(
            {
                "Operation": [
                    "Missing values fixed",
                    "Duplicate rows removed",
                    "Outliers detected and capped",
                    "Numeric columns",
                    "Categorical columns",
                    "Date columns",
                ],
                "Result": [
                    result.missing_values,
                    result.duplicates_removed,
                    result.outlier_count,
                    len(result.numeric_cols),
                    len(result.categorical_cols),
                    len(result.date_cols),
                ],
            }
        )
        st.dataframe(summary, use_container_width=True, hide_index=True)

        if result.outlier_columns:
            outliers = pd.DataFrame(
                [{"Column": column, "Outliers": count} for column, count in result.outlier_columns.items()]
            )
            fig = px.bar(outliers, x="Column", y="Outliers", color="Outliers", color_continuous_scale="Magenta")
            st.plotly_chart(style_plot(fig, height=300), use_container_width=True)

    with col2:
        st.markdown('<div class="section-title">Cleaned Data Preview</div>', unsafe_allow_html=True)
        st.dataframe(result.clean_df.head(80), use_container_width=True)
        csv = result.clean_df.to_csv(index=False).encode("utf-8")
        st.download_button("Download Cleaned CSV", csv, "cleaned_dataset.csv", "text/csv")


def predictive_tab(result: AnalysisResult) -> None:
    col1, col2 = st.columns([1.2, 0.8])

    with col1:
        st.markdown('<div class="section-title">Sales Forecasting</div>', unsafe_allow_html=True)
        if result.forecast is None:
            st.warning("Forecasting needs at least one numeric target and five observations.")
        else:
            fig = go.Figure()
            historical = result.forecast[result.forecast["Type"] == "Historical"]
            future = result.forecast[result.forecast["Type"] == "Forecast"]
            fig.add_trace(go.Scatter(x=historical["Period"], y=historical["Actual"], mode="lines+markers", name="Actual", line=dict(color="#38d5ff")))
            fig.add_trace(go.Scatter(x=result.forecast["Period"], y=result.forecast["Forecast"], mode="lines", name="Forecast", line=dict(color="#2df5a1", dash="dot")))
            fig.add_trace(go.Scatter(x=future["Period"], y=future["Forecast"], mode="markers", name="Next 7", marker=dict(color="#f7c65b", size=9)))
            st.plotly_chart(style_plot(fig), use_container_width=True)

            st.success(
                f"Target: {result.forecast_target} | Model fit R2: "
                f"{result.forecast_score:.2f}" if result.forecast_score is not None else f"Target: {result.forecast_target}"
            )

    with col2:
        st.markdown('<div class="section-title">Predictive Analytics</div>', unsafe_allow_html=True)
        importance = predictive_importance(result.clean_df, result.numeric_cols)
        if importance is None:
            st.info("Feature importance needs at least two numeric fields and enough rows.")
        else:
            fig = px.bar(importance, x="Importance", y="Feature", orientation="h", color="Importance", color_continuous_scale="Teal")
            st.plotly_chart(style_plot(fig, height=360), use_container_width=True)

        anomalies = anomaly_signal(result.clean_df, result.numeric_cols)
        if anomalies is not None:
            st.metric("Anomaly records", int((anomalies["Anomaly"] == "Anomaly").sum()))


def segmentation_tab(result: AnalysisResult) -> None:
    col1, col2 = st.columns([1.25, 0.75])
    with col1:
        st.markdown('<div class="section-title">Customer Segmentation</div>', unsafe_allow_html=True)
        if result.segments is None:
            st.warning("Segmentation needs at least two numeric columns and six valid rows.")
        else:
            fig = px.scatter(
                result.segments,
                x="PCA 1",
                y="PCA 2",
                color="Segment",
                hover_data=result.categorical_cols[:3] + result.numeric_cols[:3],
                color_discrete_sequence=["#38d5ff", "#2df5a1", "#f7c65b", "#ff4da6"],
            )
            st.plotly_chart(style_plot(fig), use_container_width=True)

    with col2:
        st.markdown('<div class="section-title">Segment Profile</div>', unsafe_allow_html=True)
        if result.segment_profile is not None:
            st.dataframe(result.segment_profile, use_container_width=True, hide_index=True)
        else:
            st.info("No segment profile available.")


def insights_tab(result: AnalysisResult) -> None:
    left, right = st.columns([0.9, 1.1])
    with left:
        st.markdown('<div class="section-title">Actionable Analytical Insights</div>', unsafe_allow_html=True)
        for insight in result.insights:
            st.markdown(f'<div class="insight">{insight}</div>', unsafe_allow_html=True)

        pdf_bytes = build_pdf_report(result)
        ppt_bytes = build_powerpoint_report(result)
        export_col1, export_col2 = st.columns(2)
        export_col1.download_button("PDF Report", pdf_bytes, "data_analytics_report.pdf", "application/pdf")
        export_col2.download_button(
            "PowerPoint Export",
            ppt_bytes,
            "data_analytics_report.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    with right:
        st.markdown('<div class="section-title">AI Chat Assistant</div>', unsafe_allow_html=True)
        if "messages" not in st.session_state:
            st.session_state.messages = [
                ("ai", "Ask me about data quality, missing values, outliers, sales forecast, or customer segments.")
            ]

        for role, message in st.session_state.messages[-8:]:
            css_class = "chat-user" if role == "user" else "chat-ai"
            st.markdown(f'<div class="{css_class}">{message}</div>', unsafe_allow_html=True)

        question = st.chat_input("Ask for a decision-ready insight")
        if question:
            st.session_state.messages.append(("user", question))
            st.session_state.messages.append(("ai", answer_question(question, result)))
            st.rerun()


def answer_question(question: str, result: AnalysisResult) -> str:
    text = question.lower()
    if "missing" in text or "clean" in text:
        return (
            f"The system handled {result.missing_values:,} missing values, removed "
            f"{result.duplicates_removed:,} duplicate rows, and capped {result.outlier_count:,} outliers."
        )
    if "quality" in text or "score" in text:
        return (
            f"The current data quality score is {result.quality['score']}/100. "
            f"Completeness is {result.quality['completeness']}%, uniqueness is {result.quality['uniqueness']}%, "
            f"and consistency is {result.quality['consistency']}%."
        )
    if "outlier" in text:
        if not result.outlier_columns:
            return "No strong IQR outliers were found in the numeric columns."
        columns = ", ".join(f"{column} ({count})" for column, count in result.outlier_columns.items())
        return f"Outliers were detected in: {columns}. These values were capped to reduce model distortion."
    if "forecast" in text or "sales" in text or "predict" in text:
        if result.forecast is None:
            return "Forecasting is unavailable because the dataset needs more numeric observations."
        next_value = result.forecast.loc[result.forecast["Type"] == "Forecast", "Forecast"].iloc[0]
        return f"The next forecasted {result.forecast_target} value is approximately {next_value:,.2f}."
    if "segment" in text or "customer" in text:
        if result.segment_profile is None:
            return "Segmentation needs at least two numeric columns and six valid rows."
        return " ".join(
            f"{row.Segment}: {int(row.Records)} records, average target {row.Average_Target:,.2f}."
            for row in result.segment_profile.itertuples()
        )
    return " ".join(result.insights[:3])


def build_pdf_report(result: AnalysisResult) -> bytes:
    buffer = BytesIO()
    document = SimpleDocTemplate(buffer, pagesize=letter)
    styles = getSampleStyleSheet()
    story: list[Any] = [
        Paragraph("AI Data Cleaning and Visualization Report", styles["Title"]),
        Spacer(1, 12),
        Paragraph(f"Data Quality Score: {result.quality['score']}/100", styles["Heading2"]),
        Spacer(1, 8),
    ]

    metrics = [
        ["Rows", f"{len(result.clean_df):,}"],
        ["Columns", f"{result.clean_df.shape[1]:,}"],
        ["Missing Fixed", f"{result.missing_values:,}"],
        ["Duplicates Removed", f"{result.duplicates_removed:,}"],
        ["Outliers", f"{result.outlier_count:,}"],
    ]
    table = Table(metrics, colWidths=[220, 180])
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#eef5ff")),
                ("TEXTCOLOR", (0, 0), (-1, -1), colors.HexColor("#182033")),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#9eb2cb")),
                ("FONTNAME", (0, 0), (-1, -1), "Helvetica-Bold"),
                ("PADDING", (0, 0), (-1, -1), 8),
            ]
        )
    )
    story.extend([table, Spacer(1, 16), Paragraph("Actionable Insights", styles["Heading2"])])

    for insight in result.insights:
        story.append(Paragraph(f"- {insight}", styles["BodyText"]))
        story.append(Spacer(1, 6))

    if result.segment_profile is not None:
        story.extend([Spacer(1, 12), Paragraph("Customer Segments", styles["Heading2"])])
        segment_rows = [["Segment", "Records", "Average Target"]] + result.segment_profile.round(2).astype(str).values.tolist()
        story.append(Table(segment_rows, colWidths=[150, 100, 150]))

    document.build(story)
    buffer.seek(0)
    return buffer.getvalue()


def build_powerpoint_report(result: AnalysisResult) -> bytes:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    add_slide(
        presentation,
        "AI Data Cleaning and Visualization System",
        [
            f"Quality score: {result.quality['score']}/100",
            f"Clean rows: {len(result.clean_df):,}",
            f"Columns profiled: {result.clean_df.shape[1]:,}",
            "Built with Python, Streamlit, scikit-learn, and Plotly",
        ],
    )
    add_slide(
        presentation,
        "Automated Data Cleaning",
        [
            f"Missing values fixed: {result.missing_values:,}",
            f"Duplicates removed: {result.duplicates_removed:,}",
            f"Outliers detected and capped: {result.outlier_count:,}",
            f"Completeness: {result.quality['completeness']}%",
        ],
    )
    add_slide(presentation, "Actionable Insights", result.insights[:5])

    predictive_points = []
    if result.forecast is not None:
        next_value = result.forecast.loc[result.forecast["Type"] == "Forecast", "Forecast"].iloc[0]
        predictive_points.append(f"Forecast target: {result.forecast_target}")
        predictive_points.append(f"Next forecast estimate: {next_value:,.2f}")
    if result.segment_profile is not None:
        predictive_points.append(f"Customer segments generated: {len(result.segment_profile)}")
    add_slide(presentation, "Predictive Analytics and Segmentation", predictive_points or ["Not enough data for predictive models."])

    buffer = BytesIO()
    presentation.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()


def add_slide(presentation: Presentation, title: str, bullets: list[str]) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(7, 17, 31)

    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.45), Inches(12), Inches(0.75))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(30)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(234, 242, 255)

    body = slide.shapes.add_textbox(Inches(0.95), Inches(1.45), Inches(11.4), Inches(5.6))
    frame = body.text_frame
    frame.word_wrap = True
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(22)
        paragraph.font.color.rgb = RGBColor(180, 203, 229)


def sidebar_data_loader() -> pd.DataFrame | None:
    st.sidebar.markdown("## Data Input")
    uploaded_file = st.sidebar.file_uploader("Upload CSV or Excel dataset", type=["csv", "xlsx", "xls"])
    use_sample = st.sidebar.toggle("Use sample sales dataset", value=uploaded_file is None)

    if uploaded_file is not None:
        try:
            return read_uploaded_file(uploaded_file)
        except Exception as exc:
            st.sidebar.error(str(exc))
            return None

    if use_sample:
        return sample_dataset()

    return None


def main() -> None:
    inject_css()
    show_hero()

    raw_df = sidebar_data_loader()
    if raw_df is None or raw_df.empty:
        st.info("Upload a dataset or enable the sample dataset to start analysis.")
        return

    result = analyze(raw_df)

    metric_cols = st.columns(5)
    with metric_cols[0]:
        metric_card("Clean Rows", f"{len(result.clean_df):,}", "ready for analysis")
    with metric_cols[1]:
        metric_card("Columns", f"{result.clean_df.shape[1]:,}", "profiled")
    with metric_cols[2]:
        metric_card("Quality", f"{result.quality['score']}/100", "AI score")
    with metric_cols[3]:
        metric_card("Missing Fixed", f"{result.missing_values:,}", "auto-imputed")
    with metric_cols[4]:
        metric_card("Outliers", f"{result.outlier_count:,}", "detected")

    tabs = st.tabs(
        [
            "Auto Dashboard",
            "Data Cleaning",
            "Predictive Analytics",
            "Customer Segmentation",
            "AI Assistant and Reports",
        ]
    )
    with tabs[0]:
        dashboard_tab(result)
    with tabs[1]:
        cleaning_tab(result)
    with tabs[2]:
        predictive_tab(result)
    with tabs[3]:
        segmentation_tab(result)
    with tabs[4]:
        insights_tab(result)


if __name__ == "__main__":
    main()
